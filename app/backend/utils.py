from pathlib import Path

import docx
import PyPDF2
from io import BytesIO
from pptx import Presentation
import pytesseract
from PIL import Image
import torch
from captum.attr import LayerGradCam
from model.transformer import TransformerClassifier, tokenizer, MAX_LENGTH


def extract_text_from_txt(content: bytes) -> str:
    return content.decode('utf-8')


def extract_text_from_docx(content: bytes) -> str:
    doc = docx.Document(BytesIO(content))
    return '\n'.join([paragraph.text for paragraph in doc.paragraphs])


def extract_text_from_pdf(content: bytes) -> str:
    pdf_file = BytesIO(content)
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    text = ''
    for page in pdf_reader.pages:
        text += page.extract_text() + '\n'
    return text


def extract_text_from_pptx(content: bytes) -> str:
    prs = Presentation(BytesIO(content))
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text.append(shape.text)

    return '\n'.join(text)


# TODO: extract text from image using Gemini
def extract_text_from_image(content: bytes) -> str:
    try:
        # Open the image using PIL
        image = Image.open(BytesIO(content))

        # Convert to RGB if necessary (for PNG with transparency)
        if image.mode in ('RGBA', 'LA'):
            background = Image.new('RGB', image.size, (255, 255, 255))
            background.paste(image, mask=image.split()[-1])
            image = background
        elif image.mode != 'RGB':
            image = image.convert('RGB')

        # Extract text using Tesseract OCR
        text = pytesseract.image_to_string(image, lang='rus+eng')

        return text.strip()
    except Exception as e:
        raise ValueError(f'Failed to process image: {str(e)}')


def analyze_text_with_gradcam(text: str) -> list[dict[str, float]]:
    # Load model and weights
    device = torch.device('cuda' if torch.cuda.is_available() else 'cpu')
    model = TransformerClassifier(
        vocab_size=tokenizer.vocab_size, d_model=256, nhead=8, num_layers=6, dim_feedforward=1024, dropout=0.1
    )
    path = (
            Path(__file__).resolve().parent.parent
            / 'model'
            / 'transformer.pth'
    )
    model.load_state_dict(torch.load(path, map_location=device))
    model.to(device)
    model.eval()

    # Tokenize input text
    encoding = tokenizer(text, max_length=MAX_LENGTH, padding='max_length', truncation=True, return_tensors='pt')
    input_ids = encoding['input_ids'].to(device)
    attention_mask = encoding['attention_mask'].to(device)

    # Initialize Grad-CAM
    grad_cam = LayerGradCam(model, model.transformer_encoder)

    # Get Grad-CAM attributions for AI-written class (0)
    attributions = grad_cam.attribute(
        input_ids,
        target=0,  # Target class (AI-written = 0, Human-written = 1)
        additional_forward_args=(attention_mask,),
    )

    # Process tokens and their scores
    tokens = tokenizer.convert_ids_to_tokens(input_ids[0])
    # Detach the tensor before converting to numpy
    token_scores = attributions[0].detach().cpu().numpy()[0]

    # Create result list with token-level scores
    results = []
    for token, score in zip(tokens, token_scores):
        if token not in ['[PAD]', '[CLS]', '[SEP]']:
            # Normalize score to probability range [0, 1]
            normalized_score = float((score - token_scores.min()) / (token_scores.max() - token_scores.min()))
            results.append(
                {
                    'token': token,
                    'ai_prob': normalized_score,  # Higher score means more likely to be AI-written
                    'is_special_token': token.startswith('[') and token.endswith(']'),
                }
            )

    return results
